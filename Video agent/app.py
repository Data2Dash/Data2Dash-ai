"""
AI Video Presenter Pro — Creative Edition (v4)
===============================================
Changes from v3:
  - Captions removed entirely
  - Fixed closure bug (same frame repeating across all slides)
  - Unique illustration per slide (seed = slide_index * 7919 + topic_hash)
  - 8 visually distinct slide themes — each with its own color palette,
    background style, layout, and typography treatment
  - Backgrounds: gradient washes, bold color blocks, split-color, paper texture,
    dark mode, blueprint, warm cream — never the same boring dot-grid
"""

# ── Imports ───────────────────────────────────────────────────────────────────
import streamlit as st
import os, re, asyncio, tempfile, base64, textwrap, requests, urllib.parse, uuid, math
from io import BytesIO
from typing import Optional, List, Tuple

import fitz
import docx as python_docx
from pptx import Presentation as PPTXPresentation
import edge_tts
from PIL import Image, ImageDraw, ImageFont, ImageFilter, ImageEnhance
import numpy as np

from moviepy import AudioFileClip, concatenate_videoclips, ImageClip

from langchain_groq import ChatGroq
from langchain_core.messages import HumanMessage
from langchain_core.prompts import ChatPromptTemplate
from pydantic import BaseModel, Field
from langchain_core.output_parsers import PydanticOutputParser

# ── Pydantic Models ───────────────────────────────────────────────────────────

class Slide(BaseModel):
    slide_type: str = Field(
        description=(
            "Layout type. Pick the best fit:\n"
            "  'title'      — Opening slide, big centered title + illustration\n"
            "  'hook'       — Large rhetorical question, centered, bold\n"
            "  'split'      — Two side-by-side comparison panels\n"
            "  'fullart'    — Illustration fills entire slide, title overlaid\n"
            "  'concept'    — Title + large text explanation + illustration\n"
            "  'timeline'   — Steps or phases listed left-to-right\n"
            "  'quote'      — Large pullquote or key statement, minimal design\n"
            "  'summary'    — Final recap slide with checkmark list\n"
            "Never use same type twice in a row."
        )
    )
    title: str = Field(description="Slide title, max 8 words. Can be empty for hook/quote.")
    hook_question: str = Field(default="", description="hook: full rhetorical question.")
    hook_highlight: str = Field(default="", description="hook: 1-3 words to yellow-highlight.")
    left_panel: str = Field(default="", description="split: left text. Wrap key terms in **asterisks**.")
    left_label: str = Field(default="", description="split: short label for left panel (1-3 words).")
    right_panel: str = Field(default="", description="split: right text. Wrap key terms in **asterisks**.")
    right_label: str = Field(default="", description="split: short label for right panel (1-3 words).")
    concept_text: str = Field(default="", description="concept: explanation text. Wrap key terms in **asterisks**.")
    timeline_steps: List[str] = Field(default=[], description="timeline: 3-5 step labels (short, 2-4 words each).")
    timeline_desc: List[str] = Field(default=[], description="timeline: description for each step (1 sentence each).")
    quote_text: str = Field(default="", description="quote: the big statement or key fact to display.")
    summary_points: List[str] = Field(default=[], description="summary: 3-5 recap points, max 12 words each.")
    image_prompt: str = Field(
        description=(
            "Image generation prompt. Be specific and creative. "
            "For light-background slides: 'hand-drawn sketch, black ink on white, educational doodle, no text, clean lines, [subject]'. "
            "For dark/colorful slides: 'flat vector illustration, colorful, [subject], no text, clean minimal style'. "
            "Make each prompt completely unique and visually distinct."
        )
    )
    script: str = Field(
        description=(
            "Spoken narration — MINIMUM 120 WORDS, natural conversational sentences. "
            "Structure: intro (15w) → explain with examples (60w) → analogy (30w) → implications (20w) → transition (15w). "
            "No markdown, no bullets. Use '...' for pauses."
        )
    )

class PresentationModel(BaseModel):
    presentation_title: str
    slides: List[Slide]

# ── Design System ─────────────────────────────────────────────────────────────
# 8 visually distinct themes, assigned round-robin per slide

THEMES = [
    # 0: Warm cream + deep navy text
    {
        "bg": [(252,248,240),(245,238,225)],   # cream gradient top→bottom
        "panel": (255,255,255),
        "accent": (30,60,120),
        "accent2": (200,80,40),
        "text": (25,22,18),
        "text2": (80,70,55),
        "style": "cream",
    },
    # 1: Bold coral/orange left block + white right
    {
        "bg": [(255,255,255),(255,255,255)],
        "panel": (245,100,60),
        "accent": (245,100,60),
        "accent2": (255,200,50),
        "text": (20,20,20),
        "text2": (255,255,255),
        "style": "bold_split",
    },
    # 2: Soft mint green throughout
    {
        "bg": [(220,245,232),(200,235,215)],
        "panel": (255,255,255),
        "accent": (30,140,90),
        "accent2": (0,100,60),
        "text": (15,50,30),
        "text2": (50,100,70),
        "style": "mint",
    },
    # 3: Dark slate — "dark mode"
    {
        "bg": [(22,28,40),(15,20,30)],
        "panel": (35,45,65),
        "accent": (100,180,255),
        "accent2": (80,220,160),
        "text": (240,240,250),
        "text2": (160,180,210),
        "style": "dark",
    },
    # 4: Lavender purple
    {
        "bg": [(240,232,255),(228,215,252)],
        "panel": (255,255,255),
        "accent": (110,60,200),
        "accent2": (200,60,180),
        "text": (40,20,80),
        "text2": (90,60,140),
        "style": "lavender",
    },
    # 5: Blueprint — blue grid paper feel
    {
        "bg": [(220,232,252),(205,220,245)],
        "panel": (255,255,255),
        "accent": (20,80,180),
        "accent2": (40,160,220),
        "text": (10,30,80),
        "text2": (40,80,150),
        "style": "blueprint",
    },
    # 6: Warm peach/rose
    {
        "bg": [(255,235,225),(252,220,208)],
        "panel": (255,255,255),
        "accent": (200,60,80),
        "accent2": (240,140,60),
        "text": (60,20,20),
        "text2": (120,50,40),
        "style": "rose",
    },
    # 7: Bright yellow editorial
    {
        "bg": [(255,245,180),(255,235,140)],
        "panel": (40,35,25),
        "accent": (40,35,25),
        "accent2": (200,40,40),
        "text": (30,25,10),
        "text2": (80,70,20),
        "style": "yellow",
    },
]

W, H = 1280, 720
YELLOW_HL = (255, 213, 0)

# ── Font Loader ───────────────────────────────────────────────────────────────

@st.cache_resource
def load_fonts():
    spec = {
        "title_xl":  [("arialbd.ttf",86), ("DejaVuSans-Bold.ttf",86)],
        "title_lg":  [("arialbd.ttf",60), ("DejaVuSans-Bold.ttf",60)],
        "title_md":  [("arialbd.ttf",42), ("DejaVuSans-Bold.ttf",42)],
        "title_sm":  [("arialbd.ttf",30), ("DejaVuSans-Bold.ttf",30)],
        "body_xl":   [("arial.ttf",  40), ("DejaVuSans.ttf",     40)],
        "body_lg":   [("arial.ttf",  32), ("DejaVuSans.ttf",     32)],
        "body_md":   [("arial.ttf",  26), ("DejaVuSans.ttf",     26)],
        "body_sm":   [("arial.ttf",  20), ("DejaVuSans.ttf",     20)],
        "label":     [("arialbd.ttf",22), ("DejaVuSans-Bold.ttf",22)],
        "brand":     [("arial.ttf",  18), ("DejaVuSans.ttf",     18)],
        "hook_xl":   [("arialbd.ttf",72), ("DejaVuSans-Bold.ttf",72)],
        "quote_xl":  [("arialbd.ttf",52), ("DejaVuSans-Bold.ttf",52)],
    }
    fonts = {}
    for key, opts in spec.items():
        for fname, size in opts:
            try:
                fonts[key] = ImageFont.truetype(fname, size)
                break
            except Exception:
                pass
        if key not in fonts:
            fonts[key] = ImageFont.load_default(size=opts[0][1])
    return fonts

# ── Background Generators ─────────────────────────────────────────────────────

def make_gradient_bg(color_top: tuple, color_bot: tuple) -> Image.Image:
    img = Image.new("RGB", (W, H))
    draw = ImageDraw.Draw(img)
    for y in range(H):
        t = y / H
        r = int(color_top[0] + t*(color_bot[0]-color_top[0]))
        g = int(color_top[1] + t*(color_bot[1]-color_top[1]))
        b = int(color_top[2] + t*(color_bot[2]-color_top[2]))
        draw.line([(0,y),(W,y)], fill=(r,g,b))
    return img

def make_blueprint_bg(accent: tuple) -> Image.Image:
    """Grid-paper style background."""
    bg_color = (218, 230, 250)
    img = Image.new("RGB", (W, H), bg_color)
    draw = ImageDraw.Draw(img)
    grid_color = (180, 205, 240)
    for x in range(0, W, 32):
        w = 2 if x % 160 == 0 else 1
        draw.line([(x,0),(x,H)], fill=grid_color, width=w)
    for y in range(0, H, 32):
        w = 2 if y % 160 == 0 else 1
        draw.line([(0,y),(W,y)], fill=grid_color, width=w)
    return img

def make_dark_bg(color_top: tuple, color_bot: tuple) -> Image.Image:
    img = make_gradient_bg(color_top, color_bot)
    # Subtle noise for texture
    arr = np.array(img, dtype=np.int16)
    rng = np.random.default_rng(42)
    noise = rng.integers(-6, 6, arr.shape, dtype=np.int16)
    arr = np.clip(arr + noise, 0, 255).astype(np.uint8)
    return Image.fromarray(arr)

def make_cream_bg() -> Image.Image:
    img = make_gradient_bg((252,248,240),(242,234,218))
    draw = ImageDraw.Draw(img)
    # Faint dot texture
    for x in range(0, W, 36):
        for y in range(0, H, 36):
            draw.ellipse([(x-1,y-1),(x+1,y+1)], fill=(220,210,190))
    return img

def make_slide_bg(theme: dict) -> Image.Image:
    style = theme["style"]
    top, bot = theme["bg"]
    if style == "dark":
        return make_dark_bg(top, bot)
    elif style == "blueprint":
        return make_blueprint_bg(theme["accent"])
    elif style == "cream":
        return make_cream_bg()
    else:
        return make_gradient_bg(top, bot)

# ── Text Helpers ──────────────────────────────────────────────────────────────

def mw(draw, text, font):
    bb = draw.textbbox((0,0), text, font=font)
    return bb[2]-bb[0], bb[3]-bb[1]

def draw_wrapped(draw, text, x, y, font, fill, max_w, spacing=10) -> int:
    words = text.split()
    lines, cur = [], []
    for w in words:
        test = " ".join(cur+[w])
        if mw(draw,test,font)[0] > max_w and cur:
            lines.append(" ".join(cur)); cur=[w]
        else:
            cur.append(w)
    if cur: lines.append(" ".join(cur))
    cy = y
    for line in lines:
        draw.text((x,cy), line, font=font, fill=fill)
        _, lh = mw(draw, line, font)
        cy += lh + spacing
    return cy

def parse_hl(text:str) -> Tuple[str, List[str]]:
    hl = re.findall(r'\*\*(.+?)\*\*', text)
    clean = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    return clean, hl

def draw_highlighted_wrapped(draw, text, x, y, font, fill, highlights, max_w, spacing=14, hl_color=YELLOW_HL) -> int:
    hl_set = {h.lower().strip() for h in highlights}
    words = text.split()
    sw, _ = mw(draw," ",font)
    lines, cur_line, cur_w = [], [], 0
    for w in words:
        ww, _ = mw(draw,w,font)
        if cur_w + ww > max_w and cur_line:
            lines.append(cur_line); cur_line=[w]; cur_w=ww+sw
        else:
            cur_line.append(w); cur_w+=ww+sw
    if cur_line: lines.append(cur_line)
    cy = y
    for line in lines:
        cx = x
        for word in line:
            cw2 = re.sub(r'[.,!?;:]','',word).lower()
            ww, wh = mw(draw,word,font)
            if cw2 in hl_set:
                draw.rectangle([(cx-5,cy-3),(cx+ww+5,cy+wh+4)], fill=hl_color)
            draw.text((cx,cy), word, font=font, fill=fill)
            cx += ww+sw
        _, lh = mw(draw," ".join(line),font)
        cy += lh+spacing
    return cy

def draw_brand(draw, fonts, theme):
    txt = "AI Presenter"
    tw, _ = mw(draw, txt, fonts["brand"])
    color = theme["text2"]
    draw.text((W-tw-24, H-30), txt, font=fonts["brand"], fill=(*color, 180))

# ── Illustration Fetcher ──────────────────────────────────────────────────────

def fetch_illustration(prompt: str, slide_idx: int, width: int = 580,
                        height: int = 460, dark_mode: bool = False) -> Image.Image:
    """
    Fetch from pollinations with a unique seed per slide index.
    Falls back to a PIL-drawn placeholder on failure.
    """
    # Use slide index in seed to guarantee unique image per slide
    seed = (slide_idx * 7919 + abs(hash(prompt[:40]))) % 999983

    style_suffix = (
        ", flat vector art, vibrant colors, no text, bold shapes"
        if dark_mode else
        ", hand-drawn sketch, black ink on white, educational doodle, no text, clean lines"
    )
    full_prompt = prompt.rstrip(".") + style_suffix

    try:
        safe = urllib.parse.quote(full_prompt)
        url  = (f"https://image.pollinations.ai/prompt/{safe}"
                f"?width={width}&height={height}&nologo=true&seed={seed}&model=flux")
        r = requests.get(url, timeout=32)
        r.raise_for_status()
        img = Image.open(BytesIO(r.content)).convert("RGBA")
        img = img.resize((width, height), Image.LANCZOS)
        arr = np.array(img.convert("L"))
        if arr.std() < 6:  # nearly blank
            raise ValueError("blank")
        return img
    except Exception:
        return _pil_fallback(prompt, width, height, slide_idx)

def _pil_fallback(title: str, w: int, h: int, idx: int) -> Image.Image:
    """Draw a clean geometric placeholder diagram."""
    import random
    rng = random.Random(idx * 31337)
    colors = [(100,160,220),(80,200,140),(220,100,80),(180,100,220),(240,180,60)]
    bg_col = (240,242,248)
    img = Image.new("RGBA", (w,h), (*bg_col,255))
    draw = ImageDraw.Draw(img)
    # Draw 3 connected boxes
    bw, bh = w//5, h//5
    gap = w//8
    total = 3*bw + 2*gap
    sx = (w-total)//2
    cy = h//2
    try: fnt = ImageFont.truetype("arial.ttf", 16)
    except: fnt = ImageFont.load_default(16)
    words = [w2 for w2 in re.split(r'\W+', title) if len(w2)>2][:3]
    if len(words) < 3: words = (words + ["Step","Process","Output"])[:3]
    for i, lbl in enumerate(words):
        c = colors[i % len(colors)]
        bx = sx + i*(bw+gap)
        by = cy - bh//2
        draw.rounded_rectangle([(bx,by),(bx+bw,by+bh)], radius=10, fill=(*c,200), outline=(80,80,80), width=2)
        tw2, th2 = draw.textbbox((0,0),lbl,font=fnt)[2:]
        draw.text((bx+(bw-tw2)//2, by+(bh-th2)//2), lbl, font=fnt, fill=(255,255,255))
        if i < 2:
            ax = bx+bw+4; ay=cy
            draw.line([(ax,ay),(ax+gap-8,ay)], fill=(80,80,80), width=3)
            draw.polygon([(ax+gap-8,ay-8),(ax+gap,ay),(ax+gap-8,ay+8)], fill=(80,80,80))
    # Scatter decorative dots
    for _ in range(12):
        rx=rng.randint(10,w-10); ry=rng.randint(10,h-10); rs=rng.randint(4,16)
        draw.ellipse([(rx-rs,ry-rs),(rx+rs,ry+rs)], outline=(160,160,180), width=2)
    return img

# ── Slide Renderers ───────────────────────────────────────────────────────────

def render_title(slide, fonts, theme, illus) -> Image.Image:
    canvas = make_slide_bg(theme).convert("RGBA")
    draw = ImageDraw.Draw(canvas)

    # Large decorative circle accent
    ac = theme["accent"]
    overlay = Image.new("RGBA", (W,H),(0,0,0,0))
    od = ImageDraw.Draw(overlay)
    od.ellipse([(W-340,-100),(W+100,340)], fill=(*ac, 30))
    od.ellipse([(-80,H-300),(280,H+80)], fill=(*ac, 20))
    canvas = Image.alpha_composite(canvas, overlay)
    draw = ImageDraw.Draw(canvas)

    # Illustration bottom-right or center
    if illus:
        iw,ih = 560,420
        ill = illus.resize((iw,ih), Image.LANCZOS)
        canvas.paste(ill, (W-iw-20, H-ih-10), ill)

    # Title left-aligned, large
    lines = textwrap.wrap(slide.title, width=20)
    ty = 80
    for line in lines:
        # Accent underline on first line
        if ty == 80:
            lw,lh = mw(draw, line, fonts["title_xl"])
            draw.rectangle([(60, ty+lh+4),(60+lw, ty+lh+10)], fill=(*ac,255))
        draw.text((60, ty), line, font=fonts["title_xl"], fill=tuple(theme["text"]))
        _, lh = mw(draw, line, fonts["title_xl"])
        ty += lh+8

    draw_brand(draw, fonts, theme)
    return canvas.convert("RGB")


def render_hook(slide, fonts, theme, illus) -> Image.Image:
    canvas = make_slide_bg(theme).convert("RGBA")
    draw = ImageDraw.Draw(canvas)

    # Huge background accent shape
    ac = theme["accent"]
    overlay = Image.new("RGBA",(W,H),(0,0,0,0))
    od = ImageDraw.Draw(overlay)
    od.ellipse([(W//2-350,H//2-350),(W//2+350,H//2+350)], fill=(*ac,15))
    canvas = Image.alpha_composite(canvas, overlay)
    draw = ImageDraw.Draw(canvas)

    if illus:
        # Small icons at corners
        for pos,sz in [((20,20),160),(W-180,20),(20,H-180),(W-180,H-180)]:
            px = pos if isinstance(pos, int) else pos
            if isinstance(pos,tuple): px,py = pos
            else: px,py = pos,sz; sz=160
            sm = illus.resize((sz,sz), Image.LANCZOS)
            canvas.paste(sm,(px,py),sm)

    question = slide.hook_question or slide.title
    hl = (slide.hook_highlight or "").strip()
    lines = textwrap.wrap(question, width=28)
    total_h = sum(mw(draw,l,fonts["hook_xl"])[1]+18 for l in lines)
    sy = (H-total_h)//2 - 10

    for line in lines:
        lw,lh = mw(draw,line,fonts["hook_xl"])
        ll = line.lower(); hll = hl.lower()
        if hl and hll in ll:
            idx = ll.index(hll)
            bef = line[:idx]; hlp = line[idx:idx+len(hl)]; aft = line[idx+len(hl):]
            bw,_ = mw(draw,bef,fonts["hook_xl"])
            hw,hh2 = mw(draw,hlp,fonts["hook_xl"])
            aw,_ = mw(draw,aft,fonts["hook_xl"])
            sx2 = (W-bw-hw-aw)//2
            if bef: draw.text((sx2,sy),bef,font=fonts["hook_xl"],fill=tuple(theme["text"]))
            draw.rectangle([(sx2+bw-8,sy-5),(sx2+bw+hw+8,sy+hh2+7)],fill=YELLOW_HL)
            draw.text((sx2+bw,sy),hlp,font=fonts["hook_xl"],fill=tuple(theme["text"]))
            if aft: draw.text((sx2+bw+hw,sy),aft,font=fonts["hook_xl"],fill=tuple(theme["text"]))
        else:
            draw.text(((W-lw)//2,sy),line,font=fonts["hook_xl"],fill=tuple(theme["text"]))
        sy += lh+18

    draw_brand(draw,fonts,theme)
    return canvas.convert("RGB")


def render_split(slide, fonts, theme, illus_l, illus_r) -> Image.Image:
    canvas = make_slide_bg(theme).convert("RGBA")
    draw = ImageDraw.Draw(canvas)
    ac = theme["accent"]; ac2 = theme["accent2"]

    margin,gap = 36,16
    cw = (W-2*margin-gap)//2
    ch = H-2*margin

    # Two pastel panel colors derived from theme
    panel_cols = [
        tuple(min(255,c+60) for c in ac),  # lighter accent
        tuple(min(255,c+60) for c in ac2),
    ]

    for si,side in enumerate(("left","right")):
        cx = margin if side=="left" else margin+cw+gap
        pcol = panel_cols[si]
        text_raw = slide.left_panel if side=="left" else slide.right_panel
        label    = slide.left_label if side=="left" else slide.right_label
        illus    = illus_l if side=="left" else illus_r
        clean,highlights = parse_hl(text_raw)

        # Card
        card = Image.new("RGBA",(cw,ch),(*pcol,240))
        canvas.paste(card,(cx,margin))

        # Top color bar
        bar_col = ac if si==0 else ac2
        canvas.paste(Image.new("RGBA",(cw,8),(*bar_col,255)),(cx,margin))

        # Label badge
        if label:
            lw2,lh2 = mw(draw,label.upper(),fonts["label"])
            draw.rectangle([(cx+20,margin+16),(cx+lw2+44,margin+16+lh2+10)],
                           fill=(*bar_col,255))
            draw.text((cx+32,margin+21),label.upper(),font=fonts["label"],fill=(255,255,255))

        # Illustration
        if illus:
            iw2,ih2 = cw-30, 240
            ill2 = illus.resize((iw2,ih2),Image.LANCZOS)
            canvas.paste(ill2,(cx+15,margin+70),ill2)

        # Text
        ty2 = margin+330 if illus else margin+80
        draw_highlighted_wrapped(draw,clean,cx+22,ty2,fonts["body_lg"],
                                 tuple(theme["text"]),highlights,cw-44,spacing=14)

    draw_brand(draw,fonts,theme)
    return canvas.convert("RGB")


def render_fullart(slide, fonts, theme, illus) -> Image.Image:
    """Illustration fills whole slide, bold title overlay."""
    if illus:
        base = illus.resize((W,H),Image.LANCZOS).convert("RGB")
    else:
        base = make_slide_bg(theme)

    canvas = base.convert("RGBA")
    # Dark gradient overlay at bottom for text legibility
    overlay = Image.new("RGBA",(W,H),(0,0,0,0))
    od = ImageDraw.Draw(overlay)
    for y in range(H-280,H):
        a = int(200*(y-(H-280))/280)
        od.line([(0,y),(W,y)],fill=(0,0,0,a))
    # Top too
    for y in range(0,120):
        a = int(140*(1-y/120))
        od.line([(0,y),(W,y)],fill=(0,0,0,a))
    canvas = Image.alpha_composite(canvas,overlay)
    draw = ImageDraw.Draw(canvas)

    # Title at bottom
    lines = textwrap.wrap(slide.title, width=32)
    ty = H-len(lines)*80-40
    for line in lines:
        lw,lh = mw(draw,line,fonts["title_lg"])
        draw.text(((W-lw)//2,ty),line,font=fonts["title_lg"],fill=(255,255,255))
        ty += lh+8

    draw_brand(draw,fonts,{"text2":(200,200,200)})
    return canvas.convert("RGB")


def render_concept(slide, fonts, theme, illus) -> Image.Image:
    canvas = make_slide_bg(theme).convert("RGBA")
    draw = ImageDraw.Draw(canvas)
    ac = theme["accent"]

    # Title
    if slide.title:
        draw.text((52,30),slide.title,font=fonts["title_md"],fill=tuple(theme["text"]))
        tw,th = mw(draw,slide.title,fonts["title_md"])
        draw.rectangle([(52,30+th+4),(52+tw,30+th+8)],fill=(*ac,255))

    # Large pastel panel
    pc = tuple(min(255,c+90) for c in ac)
    panel = Image.new("RGBA",(W-60,H-120),(*pc,200))
    canvas.paste(panel,(28,100))

    # Illustration left
    if illus:
        iw2,ih2 = (W-60)//2-20, H-160
        ill2 = illus.resize((iw2,ih2),Image.LANCZOS)
        canvas.paste(ill2,(38,110),ill2)

    # Text right
    tx = 28+(W-60)//2+14 if illus else 68
    tw2 = (W-60)//2-30 if illus else W-120
    clean,hl = parse_hl(slide.concept_text or slide.title)
    draw_highlighted_wrapped(draw,clean,tx,130,fonts["body_lg"],
                             tuple(theme["text"]),hl,tw2,spacing=18)

    draw_brand(draw,fonts,theme)
    return canvas.convert("RGB")


def render_timeline(slide, fonts, theme, illus) -> Image.Image:
    canvas = make_slide_bg(theme).convert("RGBA")
    draw = ImageDraw.Draw(canvas)
    ac = theme["accent"]; ac2 = theme["accent2"]

    if slide.title:
        draw.text((52,28),slide.title,font=fonts["title_md"],fill=tuple(theme["text"]))

    steps = slide.timeline_steps or []
    descs = slide.timeline_desc or []
    n = max(len(steps),1)
    box_w = min(200,(W-100)//n - 20)
    box_h = 160
    total_w = n*box_w + (n-1)*30
    sx = (W-total_w)//2
    cy = 260

    # Connecting line
    draw.line([(sx+box_w//2, cy+box_h//2),(sx+total_w-box_w//2, cy+box_h//2)],
              fill=(*ac,180), width=4)

    for i,step in enumerate(steps):
        bx = sx + i*(box_w+30)
        by = cy
        # Alternating accent colors
        col = ac if i%2==0 else ac2
        lighter = tuple(min(255,c+80) for c in col)
        draw.rounded_rectangle([(bx,by),(bx+box_w,by+box_h)],
                                radius=14, fill=(*lighter,240),
                                outline=(*col,255), width=3)
        # Step number circle
        draw.ellipse([(bx+box_w//2-20,by-20),(bx+box_w//2+20,by+20)],
                     fill=(*col,255))
        num_txt = str(i+1)
        nw,nh = mw(draw,num_txt,fonts["label"])
        draw.text((bx+box_w//2-nw//2,by-nh//2),num_txt,
                  font=fonts["label"],fill=(255,255,255))
        # Step label
        draw_wrapped(draw,step,bx+10,by+14,fonts["body_sm"],
                     tuple(theme["text"]),box_w-20,8)
        # Description below box
        if i < len(descs):
            draw_wrapped(draw,descs[i],bx,by+box_h+12,
                         fonts["body_sm"],tuple(theme["text2"]),box_w,6)

    # Small illustration
    if illus and n <= 4:
        iw2,ih2 = 220,220
        ill2 = illus.resize((iw2,ih2),Image.LANCZOS)
        canvas.paste(ill2,(W-iw2-20,H-ih2-20),ill2)

    draw_brand(draw,fonts,theme)
    return canvas.convert("RGB")


def render_quote(slide, fonts, theme, illus) -> Image.Image:
    canvas = make_slide_bg(theme).convert("RGBA")
    draw = ImageDraw.Draw(canvas)
    ac = theme["accent"]

    # Big decorative quotation marks
    draw.text((50,40), "\u201C", font=fonts["title_xl"],
              fill=(*ac,80))

    # Quote text centered
    qt = slide.quote_text or slide.title
    lines = textwrap.wrap(qt, width=30)
    total_h = sum(mw(draw,l,fonts["quote_xl"])[1]+16 for l in lines)
    sy = (H-total_h)//2-20
    for line in lines:
        lw,lh = mw(draw,line,fonts["quote_xl"])
        draw.text(((W-lw)//2,sy),line,font=fonts["quote_xl"],fill=tuple(theme["text"]))
        sy += lh+16

    # Accent bar below quote
    draw.rectangle([(W//2-80, sy+10),(W//2+80, sy+16)],fill=(*ac,255))

    # Small illustration bottom-right
    if illus:
        iw2,ih2 = 240,240
        ill2 = illus.resize((iw2,ih2),Image.LANCZOS)
        canvas.paste(ill2,(W-iw2-30,H-ih2-30),ill2)

    draw_brand(draw,fonts,theme)
    return canvas.convert("RGB")


def render_summary(slide, fonts, theme, illus) -> Image.Image:
    canvas = make_slide_bg(theme).convert("RGBA")
    draw = ImageDraw.Draw(canvas)
    ac = theme["accent"]; ac2 = theme["accent2"]

    title = slide.title or "Key Takeaways"
    tw,th = mw(draw,title,fonts["title_lg"])
    draw.text(((W-tw)//2,32),title,font=fonts["title_lg"],fill=tuple(theme["text"]))
    draw.rectangle([((W-tw)//2,32+th+6),((W+tw)//2,32+th+12)],fill=(*ac,255))

    points = slide.summary_points or []
    avail = H-155-50
    step = avail//max(len(points),1)
    py = 148

    for i,pt in enumerate(points):
        col = ac if i%2==0 else ac2
        # Colored bullet circle
        draw.ellipse([(52,py+4),(84,py+36)],fill=(*col,255))
        nw2,nh2 = mw(draw,str(i+1),fonts["label"])
        draw.text((68-nw2//2,py+4+(32-nh2)//2),str(i+1),
                  font=fonts["label"],fill=(255,255,255))
        draw_wrapped(draw,pt,102,py,fonts["body_md"],tuple(theme["text"]),W-260,8)
        py += step

    if illus:
        iw2,ih2 = 280,280
        ill2 = illus.resize((iw2,ih2),Image.LANCZOS)
        canvas.paste(ill2,(W-iw2-30,H-ih2-30),ill2)

    draw_brand(draw,fonts,theme)
    return canvas.convert("RGB")


def render_slide(slide, idx, total, fonts, theme, illus1, illus2=None) -> Image.Image:
    t = slide.slide_type
    if   t == "title":    return render_title(slide,fonts,theme,illus1)
    elif t == "hook":     return render_hook(slide,fonts,theme,illus1)
    elif t == "split":    return render_split(slide,fonts,theme,illus1,illus2)
    elif t == "fullart":  return render_fullart(slide,fonts,theme,illus1)
    elif t == "concept":  return render_concept(slide,fonts,theme,illus1)
    elif t == "timeline": return render_timeline(slide,fonts,theme,illus1)
    elif t == "quote":    return render_quote(slide,fonts,theme,illus1)
    elif t == "summary":  return render_summary(slide,fonts,theme,illus1)
    else:                 return render_concept(slide,fonts,theme,illus1)

# ── Script Expander ───────────────────────────────────────────────────────────

def ensure_script_length(script:str, target:int, topic:str, llm)->str:
    if len(script.split()) >= int(target*0.85): return script
    needed = target - len(script.split())
    try:
        prompt = ChatPromptTemplate.from_messages([
            ("system",
             f"Expand this narration by ~{needed} words. Topic: {topic}. "
             f"Add examples, analogies, deeper explanation. Keep natural spoken style. "
             f"Return ONLY the expanded script."),
            ("human","{script}")
        ])
        result = (prompt|llm).invoke({"script":script})
        exp = result.content.strip()
        if len(exp.split()) > len(script.split()): return exp
    except Exception: pass
    return script

# ── TTS ───────────────────────────────────────────────────────────────────────

async def _tts(text:str, path:str, voice:str):
    clean = re.sub(r'[*#_`]','',text)
    await edge_tts.Communicate(clean, voice).save(path)

def generate_audio(text:str, path:str, voice:str):
    asyncio.run(_tts(text,path,voice))

# ── Text Extractor ────────────────────────────────────────────────────────────

def extract_text(file_path:str, ext:str, api_key:str)->str:
    text=""
    if ext==".pdf":
        with fitz.open(file_path) as doc:
            for page in doc: text+=page.get_text("text")+"\n"
    elif ext==".docx":
        doc=python_docx.Document(file_path)
        for para in doc.paragraphs: text+=para.text+"\n"
    elif ext==".pptx":
        prs=PPTXPresentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape,"text"): text+=shape.text+"\n"
    elif ext in (".png",".jpg",".jpeg"):
        with open(file_path,"rb") as f: enc=base64.b64encode(f.read()).decode()
        vision=ChatGroq(groq_api_key=api_key,model_name="llama-3.2-11b-vision-preview")
        msg=HumanMessage(content=[
            {"type":"text","text":"Extract ALL text. Return ONLY the text."},
            {"type":"image_url","image_url":{"url":f"data:image/jpeg;base64,{enc}"}},
        ])
        text=vision.invoke([msg]).content
    return text.strip()

# ── Session State ─────────────────────────────────────────────────────────────

for _k,_v in {"presentation":None,"assets_dir":None,"assets_done":False,"final_video":None}.items():
    if _k not in st.session_state: st.session_state[_k]=_v

# ══════════════════════════════════════════════════════════════════════════════
# UI
# ══════════════════════════════════════════════════════════════════════════════

st.set_page_config(page_title="AI Video Presenter",layout="wide",page_icon="🎬")
st.markdown("""
<style>
[data-testid="stAppViewContainer"]{background:#F8F7F5}
[data-testid="stSidebar"]{background:#F0EEE8}
h1{color:#1C1A18!important}
.stButton>button{background:#1C1A18;color:#F8F7F5;border:none;border-radius:8px;
  padding:.55rem 1.4rem;font-weight:600;font-size:15px}
.stButton>button:hover{background:#3a3835}
div[data-testid="metric-container"]{background:#ECEAE4;border-radius:10px;
  padding:12px 16px;border:1px solid #D8D5CC}
</style>
""", unsafe_allow_html=True)

with st.sidebar:
    st.markdown("## ⚙️ Settings")
    api_key = st.text_input("Groq API Key",type="password",help="Get free at console.groq.com")
    st.divider()
    voice_choice = st.selectbox("🎙️ Voice",[
        "en-US-AndrewNeural","en-US-AvaNeural","en-US-ChristopherNeural",
        "en-GB-SoniaNeural","en-AU-WilliamNeural","en-IN-NeerjaNeural",
    ])
    st.divider()
    num_slides     = st.slider("📑 Slides",4,16,8)
    target_minutes = st.slider("⏱️ Length (min)",2,20,6)
    words_per_slide = max(120,int((target_minutes*130)//num_slides))
    st.caption(f"~{words_per_slide} words/slide → est. **{round((words_per_slide*num_slides)/130)} min**")
    st.divider()
    st.caption("PDF · DOCX · PPTX · PNG · JPG")

if not api_key:
    st.title("🎬 AI Video Presenter")
    st.info("👈 Enter your Groq API key in the sidebar.")
    st.stop()

FONTS  = load_fonts()
llm    = ChatGroq(groq_api_key=api_key,model_name="llama-3.3-70b-versatile",temperature=0.75)
parser = PydanticOutputParser(pydantic_object=PresentationModel)

st.title("🎬 AI Video Presenter")
st.markdown("*Cinematic educational videos — unique design every slide*")
c1,c2,c3,c4 = st.columns(4)
c1.metric("Themes","8 unique styles")
c2.metric("Quality","1280×720 · 24fps")
c3.metric("Art","AI illustrations")
c4.metric("Audio","TTS narration")
st.divider()

uploaded = st.file_uploader("📎 Upload Document",type=["pdf","docx","pptx","png","jpg","jpeg"])

if uploaded and st.button("🎬  Generate Video",use_container_width=True):
    st.session_state.update({"presentation":None,"assets_dir":None,
                              "assets_done":False,"final_video":None})
    ext = os.path.splitext(uploaded.name)[1].lower()

    with tempfile.TemporaryDirectory() as tmp:
        in_path=os.path.join(tmp,f"input{ext}")
        with open(in_path,"wb") as f: f.write(uploaded.getvalue())

        with st.status("🧠 Analysing…",expanded=True) as status:
            st.write("📖 Extracting text…")
            doc_text = extract_text(in_path,ext,api_key)[:30_000]

            st.write("🎨 Designing slides…")
            fmt = parser.get_format_instructions().replace("{","{{").replace("}","}}")

            sys_prompt = f"""You are a creative educational video designer. 
Produce EXACTLY {num_slides} slides. Be creative and varied — no two slides should feel the same.

━━━ SCRIPT LENGTH (TOP PRIORITY) ━━━
Every script MUST be AT LEAST {words_per_slide} words. Count carefully.
Target: {target_minutes} min × 130 wpm = {target_minutes*130} total words across {num_slides} slides.
Structure each script: intro→explain→example→analogy→implications→transition.

━━━ SLIDE TYPES (use variety, never same type twice in a row) ━━━
Slide 1:           'title'
Slide {num_slides}: 'summary' with 4-5 summary_points
Middle slides use: 'hook', 'split', 'fullart', 'concept', 'timeline', 'quote'

'hook':     hook_question (compelling full question) + hook_highlight (1-3 words for yellow box)
'split':    left_panel + right_panel (1-2 sentences each, **highlight** key words),
            left_label + right_label (short 1-3 word headers), 
'concept':  concept_text (2-3 sentences with **highlighted** key terms)
'timeline': timeline_steps (3-5 short labels) + timeline_desc (one sentence each)
'quote':    quote_text (striking statement or key fact, 1-2 sentences)
'fullart':  just title + detailed image_prompt — illustration fills entire frame

━━━ IMAGE PROMPTS ━━━
Make each prompt UNIQUE and visually distinct from the others.
Light slides: start with "hand-drawn sketch, black ink on white, educational doodle, no text, clean lines, "
Dark/colorful slides: "flat vector illustration, vibrant colors, no text, bold shapes, "
Then describe a specific visual scene relevant to the slide topic.

{fmt}"""

            prompt = ChatPromptTemplate.from_messages([
                ("system",sys_prompt),
                ("human","Document:\n\n{text}"),
            ])
            presentation=(prompt|llm|parser).invoke({"text":doc_text})

            st.write("✍️ Checking scripts…")
            for slide in presentation.slides:
                if len(slide.script.split()) < int(words_per_slide*0.85):
                    slide.script=ensure_script_length(slide.script,words_per_slide,slide.title,llm)

            st.session_state.presentation=presentation
            adir=os.path.join(tempfile.gettempdir(),f"aipres_{uuid.uuid4().hex}")
            os.makedirs(adir,exist_ok=True)
            st.session_state.assets_dir=adir
            status.update(label="✅ Ready!",state="complete")

# ── Asset Generation ──────────────────────────────────────────────────────────

if st.session_state.presentation and not st.session_state.assets_done:
    pres  = st.session_state.presentation
    adir  = st.session_state.assets_dir
    n     = len(pres.slides)

    # Pre-build all image clips (no closures — store paths and timings list)
    slide_data = []   # list of (img_path, audio_path, duration)

    with st.status("🎬 Rendering…",expanded=True) as status:
        prog = st.progress(0.0)

        for i,slide in enumerate(pres.slides):
            prog.progress(i/n, text=f"Slide {i+1}/{n} — illustration…")
            theme = THEMES[i % len(THEMES)]
            dark = theme["style"]=="dark"

            # Illustrations — unique seed per slide via fetch_illustration(slide_idx=i)
            ill1_path=os.path.join(adir,f"ill1_{i}.png")
            ill2_path=os.path.join(adir,f"ill2_{i}.png")

            if not os.path.exists(ill1_path):
                ill1=fetch_illustration(slide.image_prompt, slide_idx=i,
                                        width=640,height=520,dark_mode=dark)
                ill1.save(ill1_path)
            ill1=Image.open(ill1_path)
            ill2=None
            if slide.slide_type=="split":
                if not os.path.exists(ill2_path):
                    alt=slide.image_prompt+" alternative perspective"
                    fetch_illustration(alt,slide_idx=i+1000,
                                       width=640,height=520,dark_mode=dark).save(ill2_path)
                ill2=Image.open(ill2_path)

            # Render slide image
            prog.progress(i/n+0.25/n, text=f"Slide {i+1}/{n} — rendering…")
            img_path=os.path.join(adir,f"img_{i}.png")
            if not os.path.exists(img_path):
                render_slide(slide,i,n,FONTS,theme,ill1,ill2).save(img_path,quality=95)

            # Audio
            prog.progress(i/n+0.55/n, text=f"Slide {i+1}/{n} — voice…")
            audio_path=os.path.join(adir,f"audio_{i}.mp3")
            if not os.path.exists(audio_path):
                generate_audio(slide.script,audio_path,voice_choice)

            audio_clip = AudioFileClip(audio_path)
            duration   = audio_clip.duration + 0.4
            audio_clip.close()

            slide_data.append((img_path, audio_path, duration))
            prog.progress((i+1)/n, text=f"Slide {i+1}/{n} ✓")

        # Build clips AFTER all assets are ready — avoids closure capture bug
        prog.progress(1.0, text="Assembling video…")
        st.write("🎞️ Encoding MP4…")

        clips = []
        for img_path, audio_path, duration in slide_data:
            ac = AudioFileClip(audio_path)
            vc = ImageClip(img_path).with_duration(duration).with_audio(ac)
            clips.append(vc)

        final_mp4=os.path.join(adir,"final_presentation.mp4")
        concatenate_videoclips(clips,method="compose").write_videofile(
            final_mp4,fps=24,codec="libx264",audio_codec="aac",
            audio_fps=44100,preset="fast",threads=4,logger=None,
        )
        st.session_state.final_video=final_mp4
        st.session_state.assets_done=True
        status.update(label="🎉 Done!",state="complete")

# ── Output ────────────────────────────────────────────────────────────────────

if st.session_state.assets_done and st.session_state.final_video:
    st.success("🎉 Your video is ready!")
    st.video(st.session_state.final_video)
    dl,rst=st.columns(2)
    with dl:
        with open(st.session_state.final_video,"rb") as vf:
            st.download_button("📥 Download MP4",data=vf.read(),
                               file_name="ai_presentation.mp4",
                               mime="video/mp4",use_container_width=True)
    with rst:
        if st.button("🔄 New Presentation",use_container_width=True):
            for k in ["assets_dir","presentation","final_video"]:
                st.session_state[k]=None
            st.session_state.assets_done=False
            st.rerun()

    with st.expander("📊 Slide Details",expanded=False):
        for i,s in enumerate(st.session_state.presentation.slides):
            wc=len(s.script.split())
            st.markdown(f"**Slide {i+1} ({s.slide_type}) — {wc} words:** {s.title}")
            with st.expander(f"🎙️ Script — slide {i+1}"):
                st.write(s.script)
            st.divider()